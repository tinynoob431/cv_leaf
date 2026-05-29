import argparse
from datetime import datetime
from pathlib import Path

import pandas as pd
from docx import Document
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.shared import Inches, Pt
from pptx import Presentation
from pptx.util import Inches as PptInches
from pptx.util import Pt as PptPt


PRESET_ORDER = ["baseline", "aug", "cbam", "cbam_mix"]
PRESET_META = {
    "baseline": ("ResNet50 baseline", "普通迁移学习"),
    "aug": ("ResNet50 + 数据增强", "提升泛化"),
    "cbam": ("ResNet50 + CBAM", "引入注意力机制"),
    "cbam_mix": ("ResNet50 + CBAM + MixUp/CutMix", "最终模型"),
}


def parse_args():
    parser = argparse.ArgumentParser(description="Generate report/ppt/readme assets for submission.")
    parser.add_argument("--root", default=".", help="Project root directory.")
    return parser.parse_args()


def read_lines(path: Path):
    return [x.strip() for x in path.read_text(encoding="utf-8").splitlines() if x.strip()]


def load_dataset_stats(root: Path):
    processed = root / "processed_data"
    labels = read_lines(processed / "label_list.txt")
    class_names = []
    for line in labels:
        if "\t" in line:
            _, name = line.split("\t", 1)
        elif " " in line and line.split(" ", 1)[0].isdigit():
            _, name = line.split(" ", 1)
        else:
            name = line
        class_names.append(name.strip())

    train_lines = read_lines(processed / "train.txt")
    eval_lines = read_lines(processed / "eval.txt")

    train_count_by_class = {i: 0 for i in range(len(class_names))}
    eval_count_by_class = {i: 0 for i in range(len(class_names))}
    for line in train_lines:
        _, y = line.rsplit(" ", 1)
        train_count_by_class[int(y)] += 1
    for line in eval_lines:
        _, y = line.rsplit(" ", 1)
        eval_count_by_class[int(y)] += 1

    return {
        "class_names": class_names,
        "num_classes": len(class_names),
        "num_train": len(train_lines),
        "num_eval": len(eval_lines),
        "train_by_class": train_count_by_class,
        "eval_by_class": eval_count_by_class,
    }


def load_ablation(root: Path):
    csv_path = root / "ablation_results.csv"
    if not csv_path.exists():
        raise FileNotFoundError(f"Missing {csv_path}")
    df = pd.read_csv(csv_path)
    df = df.drop_duplicates(subset=["ablation_preset"], keep="last")
    df["order"] = df["ablation_preset"].map({k: i for i, k in enumerate(PRESET_ORDER)})
    df = df.sort_values("order")
    return df


def ensure_dirs(root: Path):
    deliver = root / "deliverables"
    report_dir = deliver / "report"
    ppt_dir = deliver / "presentation"
    report_dir.mkdir(parents=True, exist_ok=True)
    ppt_dir.mkdir(parents=True, exist_ok=True)
    return deliver, report_dir, ppt_dir


def pick_image(root: Path, preset: str, filename: str):
    p = root / "outputs_leaf_ablation" / preset / filename
    return p if p.exists() else None


def add_title(doc: Document, text: str):
    p = doc.add_paragraph()
    p.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
    run = p.add_run(text)
    run.bold = True
    run.font.size = Pt(20)


def add_subtitle(doc: Document, text: str):
    p = doc.add_paragraph()
    p.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
    run = p.add_run(text)
    run.font.size = Pt(11)


def build_report_docx(root: Path, report_dir: Path, df: pd.DataFrame, stats: dict):
    doc = Document()
    add_title(doc, "学科实践二大作业实验报告")
    add_subtitle(doc, "题目：基于 ResNet50 + CBAM 的叶片图像分类与消融实验")
    add_subtitle(doc, f"生成日期：{datetime.now().strftime('%Y-%m-%d')}")
    doc.add_paragraph("")

    doc.add_heading("1. 课题任务描述", level=1)
    doc.add_paragraph(
        "本项目围绕图像分类任务开展，目标是完成一个不少于 5 类的视觉分类系统，"
        "并进行可复现的消融实验。我们使用迁移学习的 ResNet50 作为主干网络，"
        "进一步引入 CBAM 注意力机制和 MixUp/CutMix 增强策略，比较不同设计对验证集准确率的影响。"
    )

    doc.add_heading("2. 数据集说明", level=1)
    doc.add_paragraph(
        "数据来源：PlantVillage 公开叶片数据集（公开下载地址见参考文献），本项目离线构建为 processed_data 格式。"
    )
    doc.add_paragraph(
        f"本次实验共使用 {stats['num_classes']} 类，训练集 {stats['num_train']} 张，验证集 {stats['num_eval']} 张。"
    )
    table = doc.add_table(rows=1, cols=4)
    table.style = "Table Grid"
    hdr = table.rows[0].cells
    hdr[0].text = "类别ID"
    hdr[1].text = "类别名"
    hdr[2].text = "训练样本数"
    hdr[3].text = "验证样本数"
    for i, name in enumerate(stats["class_names"]):
        cells = table.add_row().cells
        cells[0].text = str(i)
        cells[1].text = name
        cells[2].text = str(stats["train_by_class"][i])
        cells[3].text = str(stats["eval_by_class"][i])

    img = pick_image(root, "baseline", "class_distribution.png")
    if img:
        doc.add_paragraph("图 1 类别分布：")
        doc.add_picture(str(img), width=Inches(6.2))

    doc.add_heading("3. 技术路线", level=1)
    doc.add_paragraph("整体流程：数据准备 -> 数据增强 -> 模型训练 -> 验证评估 -> 错误分析 -> 消融对比。")
    doc.add_paragraph("模型结构：ResNet50 backbone + (可选) CBAM + Dropout + 线性分类头。")
    doc.add_paragraph("训练策略：先 head-only，再 finetune；当前快速复现实验采用 1 epoch smoke setting。")

    doc.add_heading("4. 消融实验设计", level=1)
    doc.add_paragraph("为满足课程要求，完成 4 组实验：")
    for key in PRESET_ORDER:
        model, note = PRESET_META[key]
        doc.add_paragraph(f"{model}：{note}", style="List Bullet")

    doc.add_heading("5. 阶段性结果与最好结果", level=1)
    rst_table = doc.add_table(rows=1, cols=4)
    rst_table.style = "Table Grid"
    rh = rst_table.rows[0].cells
    rh[0].text = "实验配置"
    rh[1].text = "验证准确率"
    rh[2].text = "验证损失"
    rh[3].text = "说明"
    for key in PRESET_ORDER:
        model, note = PRESET_META[key]
        row = rst_table.add_row().cells
        row[0].text = model
        sub = df[df["ablation_preset"] == key]
        if len(sub) > 0:
            row[1].text = f"{float(sub.iloc[0]['final_eval_acc']) * 100:.2f}%"
            row[2].text = f"{float(sub.iloc[0]['final_eval_loss']):.4f}"
        else:
            row[1].text = "未运行"
            row[2].text = "未运行"
        row[3].text = note

    best_row = df.loc[df["final_eval_acc"].astype(float).idxmax()]
    best_name, _ = PRESET_META[best_row["ablation_preset"]]
    doc.add_paragraph(
        f"最好结果：{best_name}，验证准确率 {float(best_row['final_eval_acc']) * 100:.2f}%，"
        f"验证损失 {float(best_row['final_eval_loss']):.4f}。"
    )

    for preset, pic_name, caption in [
        ("baseline", "accuracy_curve.png", "图 2 baseline 训练曲线"),
        ("cbam", "accuracy_curve.png", "图 3 CBAM 训练曲线"),
        ("cbam_mix", "confusion_matrix.png", "图 4 CBAM+MixUp/CutMix 混淆矩阵"),
    ]:
        p = pick_image(root, preset, pic_name)
        if p:
            doc.add_paragraph(caption)
            doc.add_picture(str(p), width=Inches(6.0))

    doc.add_heading("6. 结果分析", level=1)
    doc.add_paragraph(
        "在当前 CPU + 单 epoch 快速设置下，baseline 与轻量增强配置表现相近，"
        "CBAM 与 CBAM+MixUp/CutMix 未体现优势，主要原因是训练轮次不足、"
        "注意力模块与混合增强策略对优化稳定性要求更高。"
    )
    doc.add_paragraph(
        "后续可通过增加 epoch、使用 GPU、调低学习率、提高 warmup 轮数来更公平评估注意力机制收益。"
    )

    doc.add_heading("7. 课题总结", level=1)
    doc.add_paragraph(
        "项目已完成从公开数据下载、预处理、模型训练、消融对比、可视化分析到提交材料生成的完整流程。"
        "代码支持自动准备数据集和一键切换消融配置，具备较好的复现性。"
    )

    doc.add_heading("8. 小组分工（可按实际修改）", level=1)
    div_table = doc.add_table(rows=4, cols=3)
    div_table.style = "Table Grid"
    div_table.rows[0].cells[0].text = "成员"
    div_table.rows[0].cells[1].text = "职责"
    div_table.rows[0].cells[2].text = "工作说明"
    div_table.rows[1].cells[0].text = "组长"
    div_table.rows[1].cells[1].text = "总体设计与统筹"
    div_table.rows[1].cells[2].text = "把控选题、实验进度与最终答辩"
    div_table.rows[2].cells[0].text = "组员1"
    div_table.rows[2].cells[1].text = "数据与训练"
    div_table.rows[2].cells[2].text = "数据准备、模型训练与调参"
    div_table.rows[3].cells[0].text = "组员2"
    div_table.rows[3].cells[1].text = "报告与展示"
    div_table.rows[3].cells[2].text = "实验报告、PPT 与路演讲解"

    doc.add_heading("9. 参考文献与数据来源", level=1)
    doc.add_paragraph("[1] He et al., Deep Residual Learning for Image Recognition, arXiv:1512.03385.")
    doc.add_paragraph("[2] Woo et al., CBAM: Convolutional Block Attention Module, arXiv:1807.06521.")
    doc.add_paragraph("[3] PlantVillage Dataset.")
    doc.add_paragraph("[4] Download URL: https://github.com/spMohanty/PlantVillage-Dataset")

    out_path = report_dir / "实验报告_叶片分类_CBAM.docx"
    doc.save(out_path)
    return out_path


def build_report_markdown(root: Path, report_dir: Path, df: pd.DataFrame, stats: dict):
    lines = []
    lines.append("# 学科实践二大作业实验报告（简版）")
    lines.append("")
    lines.append("## 1. 课题任务描述")
    lines.append("完成一个不少于 5 类的图像分类任务，并开展 ResNet50/CBAM/MixUp-CutMix 的消融实验。")
    lines.append("")
    lines.append("## 2. 数据集")
    lines.append(f"- 类别数：{stats['num_classes']}")
    lines.append(f"- 训练集：{stats['num_train']}")
    lines.append(f"- 验证集：{stats['num_eval']}")
    lines.append(f"- 类别：{', '.join(stats['class_names'])}")
    lines.append("")
    lines.append("## 3. 消融结果")
    lines.append("| 模型 | 验证准确率 | 验证损失 | 说明 |")
    lines.append("|---|---:|---:|---|")
    for key in PRESET_ORDER:
        model, note = PRESET_META[key]
        sub = df[df["ablation_preset"] == key]
        if len(sub) > 0:
            acc = float(sub.iloc[0]["final_eval_acc"]) * 100
            loss = float(sub.iloc[0]["final_eval_loss"])
            lines.append(f"| {model} | {acc:.2f}% | {loss:.4f} | {note} |")
        else:
            lines.append(f"| {model} | 未运行 | 未运行 | {note} |")
    lines.append("")
    lines.append("## 4. 参考链接")
    lines.append("- ResNet: https://arxiv.org/abs/1512.03385")
    lines.append("- CBAM: https://arxiv.org/abs/1807.06521")
    lines.append("- PlantVillage: https://github.com/spMohanty/PlantVillage-Dataset")
    lines.append("- 下载地址: https://github.com/spMohanty/PlantVillage-Dataset/archive/refs/heads/master.zip")
    md_path = report_dir / "实验报告_简版.md"
    md_path.write_text("\n".join(lines), encoding="utf-8")
    return md_path


def set_textbox_text(tb, text, size=20, bold=False):
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = PptPt(size)
    run.font.bold = bold


def add_title_and_bullets(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            p = body.paragraphs[0]
        else:
            p = body.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = PptPt(20)
    return slide


def build_ppt(root: Path, ppt_dir: Path, df: pd.DataFrame, stats: dict):
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)

    # Slide 1
    s1 = prs.slides.add_slide(prs.slide_layouts[0])
    s1.shapes.title.text = "学科实践二大作业"
    s1.placeholders[1].text = (
        "基于 ResNet50 + CBAM 的叶片图像分类\n"
        f"自动生成时间：{datetime.now().strftime('%Y-%m-%d')}"
    )

    # Slide 2
    add_title_and_bullets(prs, "课题任务描述", [
        "完成一个不少于 5 类的图像分类任务",
        "实现可复现训练流程（数据、训练、评估、可视化）",
        "开展 4 组消融：baseline / +增强 / +CBAM / +CBAM+MixUp-CutMix",
    ])

    # Slide 3
    s3 = add_title_and_bullets(prs, "数据集说明", [
        "数据源：PlantVillage（公开叶片数据集）",
        f"类别数：{stats['num_classes']}（{', '.join(stats['class_names'])}）",
        f"训练集：{stats['num_train']}，验证集：{stats['num_eval']}",
        "图像特点：叶片病害特征明显、纹理与斑点细节重要",
    ])
    dist_img = pick_image(root, "baseline", "class_distribution.png")
    if dist_img:
        s3.shapes.add_picture(str(dist_img), PptInches(7.0), PptInches(1.6), width=PptInches(5.9))

    # Slide 4
    add_title_and_bullets(prs, "技术路线", [
        "预处理：统一尺寸、归一化、（可选）随机裁剪/翻转/颜色扰动",
        "模型：ResNet50 backbone + Dropout + 分类头",
        "注意力：在特征层引入 CBAM（通道+空间注意力）",
        "训练：迁移学习 + AdamW + Cosine 学习率衰减",
    ])

    # Slide 5
    add_title_and_bullets(prs, "消融实验设计", [
        "Exp1: ResNet50 baseline",
        "Exp2: ResNet50 + 数据增强",
        "Exp3: ResNet50 + CBAM",
        "Exp4: ResNet50 + CBAM + MixUp/CutMix",
    ])

    # Slide 6 results table
    s6 = prs.slides.add_slide(prs.slide_layouts[5])
    s6.shapes.title.text = "实验结果对比"
    rows, cols = 5, 4
    table = s6.shapes.add_table(rows, cols, PptInches(0.6), PptInches(1.4), PptInches(12.1), PptInches(3.5)).table
    table.cell(0, 0).text = "模型"
    table.cell(0, 1).text = "验证准确率"
    table.cell(0, 2).text = "验证损失"
    table.cell(0, 3).text = "说明"
    for ridx, key in enumerate(PRESET_ORDER, start=1):
        model, note = PRESET_META[key]
        table.cell(ridx, 0).text = model
        sub = df[df["ablation_preset"] == key]
        if len(sub) > 0:
            table.cell(ridx, 1).text = f"{float(sub.iloc[0]['final_eval_acc']) * 100:.2f}%"
            table.cell(ridx, 2).text = f"{float(sub.iloc[0]['final_eval_loss']):.4f}"
        else:
            table.cell(ridx, 1).text = "未运行"
            table.cell(ridx, 2).text = "未运行"
        table.cell(ridx, 3).text = note

    # Slide 7 visuals
    s7 = prs.slides.add_slide(prs.slide_layouts[5])
    s7.shapes.title.text = "阶段性结果截图"
    p1 = pick_image(root, "cbam", "accuracy_curve.png")
    p2 = pick_image(root, "cbam_mix", "confusion_matrix.png")
    if p1:
        s7.shapes.add_picture(str(p1), PptInches(0.6), PptInches(1.3), width=PptInches(6.1))
    if p2:
        s7.shapes.add_picture(str(p2), PptInches(6.8), PptInches(1.3), width=PptInches(6.1))

    # Slide 8 analysis
    best_row = df.loc[df["final_eval_acc"].astype(float).idxmax()]
    best_name = PRESET_META[best_row["ablation_preset"]][0]
    add_title_and_bullets(prs, "结果分析", [
        f"当前快速设置下最好模型：{best_name}（{float(best_row['final_eval_acc']) * 100:.2f}%）",
        "CBAM 与 MixUp/CutMix 对训练轮次和学习率较敏感",
        "在 CPU + 1 epoch 条件下，复杂策略优势未完全体现",
        "后续建议在 GPU 上增加 epoch 做正式版对比",
    ])

    # Slide 9 summary
    add_title_and_bullets(prs, "课题总结", [
        "已完成完整工程流程：数据 -> 训练 -> 评估 -> 可视化 -> 消融",
        "代码支持自动下载数据和一键切换实验配置",
        "提交材料已自动生成，可直接按课程要求打包提交",
    ])

    # Slide 10 division
    add_title_and_bullets(prs, "小组分工（按实际替换）", [
        "组长：方案设计、进度管理、答辩统筹",
        "组员1：数据准备、训练调参、实验复现",
        "组员2：报告撰写、PPT制作、路演讲解",
    ])

    out = ppt_dir / "课程大作业汇报_叶片分类_CBAM.pptx"
    prs.save(out)
    return out


def build_submission_checklist(root: Path, deliver_dir: Path):
    checklist = [
        "# 提交材料清单",
        "",
        "根据《项目要求.pptx》，本项目已准备以下提交物：",
        "",
        "1. 实验报告（Word）",
        "2. 源码（.py）",
        "3. 数据集（processed_data + 原始下载包）",
        "4. 汇报PPT（5分钟）",
        "",
        "## 对应路径",
        f"- 报告：{(deliver_dir / 'report' / '实验报告_叶片分类_CBAM.docx').as_posix()}",
        f"- 报告简版：{(deliver_dir / 'report' / '实验报告_简版.md').as_posix()}",
        f"- 汇报PPT：{(deliver_dir / 'presentation' / '课程大作业汇报_叶片分类_CBAM.pptx').as_posix()}",
        "- 源码：leaf_classification_cbam.py, prepare_dataset.py",
        "- 数据集：processed_data/, datasets/plantvillage_master.zip, datasets/PlantVillage-Dataset-master/",
        "",
        "## 打包命名（按课程要求）",
        "文件名格式：## _组长姓名_组员1姓名_组员2姓名.zip",
        "例如：03_张三_李四_王五.zip",
    ]
    path = deliver_dir / "提交材料清单.md"
    path.write_text("\n".join(checklist), encoding="utf-8")
    return path


def main():
    args = parse_args()
    root = Path(args.root).resolve()
    _, report_dir, ppt_dir = ensure_dirs(root)

    df = load_ablation(root)
    stats = load_dataset_stats(root)

    report_doc = build_report_docx(root, report_dir, df, stats)
    report_md = build_report_markdown(root, report_dir, df, stats)
    ppt_path = build_ppt(root, ppt_dir, df, stats)
    checklist = build_submission_checklist(root, root / "deliverables")

    print("Generated:")
    print(report_doc)
    print(report_md)
    print(ppt_path)
    print(checklist)


if __name__ == "__main__":
    main()
